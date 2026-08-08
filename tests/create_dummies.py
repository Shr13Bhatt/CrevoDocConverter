import docx
from pptx import Presentation
from PIL import Image
import os

def main():
    os.makedirs('tests', exist_ok=True)
    
    # 1. Create a dummy docx
    doc = docx.Document()
    doc.add_heading('CrevoDoc Test Document', 0)
    doc.add_paragraph('This is a sample Word document generated dynamically to test the Word-to-PDF offline converter.')
    doc.add_heading('Sub-section', 1)
    doc.add_paragraph('Testing structural preservation, layout formatting, fonts and alignment.')
    doc.save('tests/dummy.docx')
    print("[OK] Created tests/dummy.docx")

    # 2. Create a dummy pptx
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
    txBox = slide.shapes.add_textbox(100000, 100000, 5000000, 1000000)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "CrevoDoc Test Presentation"
    p.font.size = docx.shared.Pt(24)
    prs.save('tests/dummy.pptx')
    print("[OK] Created tests/dummy.pptx")

    # 3. Create dummy images
    img_red = Image.new('RGB', (300, 300), color = 'red')
    img_red.save('tests/dummy.jpg')
    img_blue = Image.new('RGB', (300, 300), color = 'blue')
    img_blue.save('tests/dummy.png')
    print("[OK] Created tests/dummy.jpg and tests/dummy.png")

if __name__ == '__main__':
    main()
