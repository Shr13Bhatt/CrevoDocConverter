import sys
import os
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def main():
    if len(sys.argv) < 3:
        print("Usage: python pdf_to_pptx.py <input_pdf> <output_pptx>")
        sys.exit(1)
        
    pdf_path = sys.argv[1]
    pptx_path = sys.argv[2]
    
    print(f"Converting PDF {pdf_path} to PPTX {pptx_path}...")
    
    try:
        doc = fitz.open(pdf_path)
        prs = Presentation()
        
        # Remove the default slides
        for i in range(len(prs.slides)-1, -1, -1):
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[i]
            
        # PDF point = 1/72 inch
        # PPTX EMU = 1/914400 inch
        # 1 point = 12700 EMUs
        PT_TO_EMU = 12700
        
        # Define blank layout
        blank_layout = prs.slide_layouts[6]
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            rect = page.rect
            width_emu = int(rect.width * PT_TO_EMU)
            height_emu = int(rect.height * PT_TO_EMU)
            
            # Set presentation size based on the first slide
            if page_num == 0:
                prs.slide_width = width_emu
                prs.slide_height = height_emu
                
            slide = prs.slides.add_slide(blank_layout)
            
            # 1. Extract and Add Images
            # We track image xrefs to extract them
            image_list = page.get_images()
            for img_idx, img_info in enumerate(image_list):
                xref = img_info[0]
                try:
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    
                    # Get all rects where this image is displayed on the page
                    rects = page.get_image_rects(xref)
                    for rect_idx, r in enumerate(rects):
                        temp_img_path = f"temp_img_{page_num}_{xref}_{rect_idx}.{image_ext}"
                        with open(temp_img_path, "wb") as f:
                            f.write(image_bytes)
                            
                        # Calculate positions in EMUs
                        left = int(r.x0 * PT_TO_EMU)
                        top = int(r.y0 * PT_TO_EMU)
                        width = int(r.width * PT_TO_EMU)
                        height = int(r.height * PT_TO_EMU)
                        
                        # Add to slide
                        slide.shapes.add_picture(temp_img_path, left, top, width, height)
                        
                        # Clean up temp file
                        try:
                            os.remove(temp_img_path)
                        except:
                            pass
                except Exception as img_err:
                    print(f"Warning: Could not extract image {xref} on page {page_num}: {img_err}", file=sys.stderr)
            
            # 2. Extract and Add Text Blocks
            # Extract text blocks with layout info
            blocks = page.get_text("blocks")
            for b in blocks:
                # b = (x0, y0, x1, y1, "text", block_no, block_type)
                x0, y0, x1, y1, text, block_no, block_type = b
                
                # Check if it contains text
                text = text.strip()
                if not text or block_type != 0:
                    continue
                    
                # Calculate coordinates in EMUs
                left = int(x0 * PT_TO_EMU)
                top = int(y0 * PT_TO_EMU)
                width = int((x1 - x0) * PT_TO_EMU)
                height = int((y1 - y0) * PT_TO_EMU)
                
                # Create textbox
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True
                
                # Split paragraphs and insert
                paragraphs = text.split('\n')
                for p_idx, para_text in enumerate(paragraphs):
                    if p_idx == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = para_text
                    
                    # Basic style defaults (custom styling can be added)
                    p.font.name = 'Arial'
                    p.font.size = Pt(14)
                    
        prs.save(pptx_path)
        doc.close()
        print("Success")
        
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(2)

if __name__ == '__main__':
    main()
