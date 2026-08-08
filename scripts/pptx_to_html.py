import sys
import os
import html
import uuid
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def get_color_hex(color):
    try:
        if color.type == 1:  # RGBColor
            return f"#{color.rgb[0]:02x}{color.rgb[1]:02x}{color.rgb[2]:02x}"
    except:
        pass
    return None

def main():
    if len(sys.argv) < 4:
        print("Usage: python pptx_to_html.py <input_pptx> <output_html> <assets_dir>")
        sys.exit(1)
        
    pptx_path = sys.argv[1]
    html_path = sys.argv[2]
    assets_dir = sys.argv[3]
    
    os.makedirs(assets_dir, exist_ok=True)
    
    try:
        prs = Presentation(pptx_path)
        
        # Convert EMUs to Inches for absolute layout in CSS
        EMU_TO_INCH = 914400.0
        slide_w_in = prs.slide_width / EMU_TO_INCH
        slide_h_in = prs.slide_height / EMU_TO_INCH
        
        html_content = [
            "<!DOCTYPE html>",
            "<html>",
            "<head>",
            "<meta charset='utf-8'>",
            "<style>",
            "  @page { size: landscape; margin: 0; }",
            "  body { margin: 0; padding: 0; font-family: 'Helvetica Neue', Arial, sans-serif; background-color: #eaeaea; }",
            "  .slide {",
            f"    width: {slide_w_in}in;",
            f"    height: {slide_h_in}in;",
            "    position: relative;",
            "    background-color: #ffffff;",
            "    overflow: hidden;",
            "    page-break-after: always;",
            "    box-sizing: border-box;",
            "    margin: 0 auto;",
            "  }",
            "  .text-box {",
            "    box-sizing: border-box;",
            "    word-wrap: break-word;",
            "    white-space: pre-wrap;",
            "  }",
            "  .table-box {",
            "    border-collapse: collapse;",
            "    width: 100%;",
            "    height: 100%;",
            "  }",
            "  .table-box td {",
            "    border: 1px solid #cccccc;",
            "    padding: 4px;",
            "  }",
            "</style>",
            "</head>",
            "<body>"
        ]
        
        for slide_idx, slide in enumerate(prs.slides):
            html_content.append(f"  <div class='slide' id='slide-{slide_idx}'>")
            
            # Check for background
            # In PPTX, background can be a solid color or fill
            try:
                bg = slide.background
                fill = bg.fill
                if fill.type == 1:  # Solid color
                    color_hex = get_color_hex(fill.fore_color)
                    if color_hex:
                        html_content.append(f"    <div style='position:absolute; top:0; left:0; width:100%; height:100%; background-color:{color_hex}; z-index:-1;'></div>")
            except:
                pass
                
            # Process shapes
            for shape_idx, shape in enumerate(slide.shapes):
                left_in = shape.left / EMU_TO_INCH if shape.left else 0.0
                top_in = shape.top / EMU_TO_INCH if shape.top else 0.0
                width_in = shape.width / EMU_TO_INCH if shape.width else 0.0
                height_in = shape.height / EMU_TO_INCH if shape.height else 0.0
                
                # Check for picture
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        image_ext = image.ext
                        img_filename = f"slide_{slide_idx}_img_{shape_idx}.{image_ext}"
                        img_path = os.path.join(assets_dir, img_filename)
                        
                        with open(img_path, "wb") as f:
                            f.write(image.blob)
                            
                        # Use a relative path from the HTML location (which will be in assets' parent dir)
                        img_src_path = os.path.join(os.path.basename(assets_dir), img_filename).replace("\\", "/")
                        html_content.append(
                            f"    <img src='{img_src_path}' style='position: absolute; "
                            f"left: {left_in}in; top: {top_in}in; width: {width_in}in; height: {height_in}in; object-fit: contain;' />"
                        )
                    except Exception as pic_err:
                        print(f"Error extracting image: {pic_err}", file=sys.stderr)
                
                # Check for table
                elif shape.has_table:
                    try:
                        tbl = shape.table
                        html_content.append(
                            f"    <div style='position: absolute; left: {left_in}in; top: {top_in}in; "
                            f"width: {width_in}in; height: {height_in}in;'>"
                        )
                        html_content.append("      <table class='table-box'>")
                        for r_idx in range(len(tbl.rows)):
                            html_content.append("        <tr>")
                            for c_idx in range(len(tbl.columns)):
                                cell = tbl.cell(r_idx, c_idx)
                                text_content = []
                                for paragraph in cell.text_frame.paragraphs:
                                    para_text = html.escape(paragraph.text)
                                    text_content.append(f"<p style='margin:2px 0;'>{para_text}</p>")
                                cell_html = "".join(text_content)
                                html_content.append(f"          <td>{cell_html}</td>")
                            html_content.append("        </tr>")
                        html_content.append("      </table>")
                        html_content.append("    </div>")
                    except Exception as tbl_err:
                        print(f"Error extracting table: {tbl_err}", file=sys.stderr)
                
                # Check for text box
                elif shape.has_text_frame:
                    try:
                        tf = shape.text_frame
                        # Only add div if there is actual text content
                        if tf.text.strip():
                            html_content.append(
                                f"    <div class='text-box' style='position: absolute; left: {left_in}in; top: {top_in}in; "
                                f"width: {width_in}in; height: {height_in}in;'>"
                            )
                            
                            for para in tf.paragraphs:
                                para_align = "left"
                                if para.alignment == 1:
                                    para_align = "center"
                                elif para.alignment == 2:
                                    para_align = "right"
                                elif para.alignment == 3:
                                    para_align = "justify"
                                    
                                margin_top = para.space_before.pt if para.space_before else 0
                                margin_bottom = para.space_after.pt if para.space_after else 0
                                
                                html_content.append(f"      <p style='text-align: {para_align}; margin-top: {margin_top}pt; margin-bottom: {margin_bottom}pt;'>")
                                
                                for run in para.runs:
                                    run_text = html.escape(run.text)
                                    font = run.font
                                    styles = []
                                    
                                    if font.name:
                                        styles.append(f"font-family: '{font.name}';")
                                    if font.size:
                                        styles.append(f"font-size: {font.size.pt}pt;")
                                    if font.bold:
                                        styles.append("font-weight: bold;")
                                    if font.italic:
                                        styles.append("font-style: italic;")
                                    if font.underline:
                                        styles.append("text-decoration: underline;")
                                    
                                    color_hex = get_color_hex(font.color)
                                    if color_hex:
                                        styles.append(f"color: {color_hex};")
                                        
                                    style_str = " ".join(styles)
                                    html_content.append(f"<span style='{style_str}'>{run_text}</span>")
                                    
                                html_content.append("      </p>")
                            html_content.append("    </div>")
                    except Exception as text_err:
                        print(f"Error extracting text box: {text_err}", file=sys.stderr)
                        
            html_content.append("  </div>")
            
        html_content.append("</body>")
        html_content.append("</html>")
        
        with open(html_path, "w", encoding="utf-8") as f:
            f.write("\n".join(html_content))
            
        print("Success")
        
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(2)

if __name__ == '__main__':
    main()
